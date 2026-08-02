from __future__ import annotations

import json
import subprocess
from pathlib import Path
from statistics import fmean
from typing import Any


def execute(operation: str, payload: dict[str, Any], context: dict[str, Any]) -> dict[str, Any]:
    handlers = {
        "browser.explore": _browser,
        "document.create": _document,
        "image.process": _image,
        "code.sandbox": _code,
        "data.advanced": _data,
        "social.connector": _social,
    }
    handler = handlers.get(context["skill_name"])
    if handler is None:
        raise ValueError(f"Unknown foundation skill: {context['skill_name']}")
    return handler(operation, payload, Path(context["output_dir"]))


def _browser(operation: str, payload: dict[str, Any], output_dir: Path) -> dict[str, Any]:
    if operation not in {"inspect", "screenshot"}:
        raise ValueError("Browser skill supports inspect and screenshot.")
    url = str(payload.get("url", ""))
    if not url.startswith(("https://", "http://")):
        raise ValueError("An HTTP(S) URL is required.")
    social_hosts = ("x.com", "instagram.com", "facebook.com", "tiktok.com", "linkedin.com")
    if any(host in url.lower() for host in social_hosts) and not payload.get(
        "platform_automation_allowed"
    ):
        return {
            "blocked": True,
            "reason": (
                "Platform automation permission is not confirmed; use an official API or stop."
            ),
        }
    from playwright.sync_api import sync_playwright

    with sync_playwright() as playwright:
        browser = playwright.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": 1280, "height": 800})
        page.goto(url, wait_until="domcontentloaded", timeout=30000)
        result = {
            "url": page.url,
            "title": page.title(),
            "text": page.locator("body").inner_text()[:20000],
            "links": page.locator("a").evaluate_all(
                "els => els.slice(0, 100).map(a => ({text: a.innerText, href: a.href}))"
            ),
        }
        if operation == "screenshot":
            target = output_dir / "page.png"
            page.screenshot(path=str(target), full_page=True)
            result["screenshot"] = target.name
        browser.close()
    return result


def _document(operation: str, payload: dict[str, Any], output_dir: Path) -> dict[str, Any]:
    if operation != "create":
        raise ValueError("Document skill supports create.")
    document_type = str(payload.get("type", "docx")).lower()
    title = str(payload.get("title", "ECK Document"))
    sections = payload.get("sections", [])
    if document_type == "docx":
        from docx import Document

        document = Document()
        document.add_heading(title, 0)
        for section in sections:
            document.add_heading(str(section.get("heading", "")), level=1)
            document.add_paragraph(str(section.get("content", "")))
        target = output_dir / "document.docx"
        document.save(target)
    elif document_type == "pdf":
        from reportlab.pdfgen.canvas import Canvas

        target = output_dir / "document.pdf"
        canvas = Canvas(str(target))
        canvas.drawString(72, 780, title[:90])
        y = 750
        for section in sections:
            canvas.drawString(72, y, str(section.get("heading", ""))[:90])
            y -= 22
            canvas.drawString(72, y, str(section.get("content", ""))[:100])
            y -= 30
        canvas.save()
    elif document_type == "pptx":
        from pptx import Presentation

        presentation = Presentation()
        for section in sections or [{"heading": title, "content": ""}]:
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = str(section.get("heading", title))
            slide.placeholders[1].text = str(section.get("content", ""))
        target = output_dir / "presentation.pptx"
        presentation.save(target)
    elif document_type == "xlsx":
        from openpyxl import Workbook

        workbook = Workbook()
        sheet = workbook.active
        sheet.title = title[:31]
        rows = payload.get("rows", [])
        for row in rows:
            sheet.append(list(row))
        target = output_dir / "workbook.xlsx"
        workbook.save(target)
    else:
        raise ValueError("Supported document types: docx, pdf, pptx, xlsx.")
    return {"artifact": target.name, "type": document_type, "bytes": target.stat().st_size}


def _image(operation: str, payload: dict[str, Any], output_dir: Path) -> dict[str, Any]:
    from PIL import Image, ImageDraw

    if operation != "create":
        raise ValueError("Image skill supports create.")
    width = min(4096, max(64, int(payload.get("width", 1024))))
    height = min(4096, max(64, int(payload.get("height", 1024))))
    image = Image.new("RGB", (width, height), str(payload.get("background", "#081426")))
    draw = ImageDraw.Draw(image)
    draw.text((40, 40), str(payload.get("text", "ECK"))[:500], fill="#62f2c5")
    target = output_dir / "image.png"
    image.save(target)
    return {"artifact": target.name, "width": width, "height": height}


def _code(operation: str, payload: dict[str, Any], output_dir: Path) -> dict[str, Any]:
    if operation != "test":
        raise ValueError("Code sandbox supports test.")
    project = output_dir / "project"
    project.mkdir()
    for relative, content in dict(payload.get("files", {})).items():
        target = (project / relative).resolve()
        if project.resolve() not in target.parents:
            raise ValueError("Project file escapes the output directory.")
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(str(content), encoding="utf-8")
    command = [str(item) for item in payload.get("command", ["python", "-m", "pytest", "-q"])]
    allowed = {"python", "python3", "pytest", "node", "npm"}
    if not command or Path(command[0]).name not in allowed:
        raise ValueError("The requested test command is not allowlisted.")
    result = subprocess.run(
        command,
        cwd=project,
        capture_output=True,
        text=True,
        timeout=180,
    )
    return {
        "passed": result.returncode == 0,
        "returncode": result.returncode,
        "output": (result.stdout + result.stderr)[-12000:],
        "files": sorted(
            str(path.relative_to(project))
            for path in project.rglob("*")
            if path.is_file()
        ),
    }


def _data(operation: str, payload: dict[str, Any], output_dir: Path) -> dict[str, Any]:
    if operation != "analyze":
        raise ValueError("Advanced data skill supports analyze.")
    records = payload.get("records", [])
    numeric: dict[str, list[float]] = {}
    for record in records:
        for key, value in dict(record).items():
            if isinstance(value, (int, float)) and not isinstance(value, bool):
                numeric.setdefault(str(key), []).append(float(value))
    summary = {
        key: {"count": len(values), "min": min(values), "max": max(values), "mean": fmean(values)}
        for key, values in numeric.items()
    }
    target = output_dir / "analysis.json"
    target.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    return {"rows": len(records), "summary": summary, "artifact": target.name}


def _social(operation: str, payload: dict[str, Any], output_dir: Path) -> dict[str, Any]:
    if operation == "policy_check":
        allowed = bool(payload.get("platform_automation_allowed"))
        official_api = bool(payload.get("official_api_available"))
        return {
            "allowed": allowed,
            "official_api_preferred": official_api,
            "next_action": (
                "Use the official API adapter."
                if allowed and official_api
                else "Stop platform automation and choose another legal platform."
            ),
        }
    if not payload.get("platform_automation_allowed"):
        return {"blocked": True, "reason": "Platform rules do not permit this automation."}
    if not payload.get("official_api_configured"):
        return {"blocked": True, "reason": "An official account adapter must be configured."}
    if not payload.get("ai_disclosure_present"):
        return {"blocked": True, "reason": "Public AI/ECK disclosure is required."}
    return {
        "ready": True,
        "operation": operation,
        "detail": "Policy gate passed; the platform-specific official API adapter must execute it.",
    }
